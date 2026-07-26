"""
build_deck.py — generates the Round 1 submission deck.

Regenerate after changing any figure:  python deploy/build_deck.py
Output: docs/Autonomous-Threat-Hunter.pptx

Constrained to 7 slides: InnovaHack Round 1 guidelines require 6-7. No template
was circulated, so the structure is ours, ordered to cover all seven evaluation
criteria in weight order.

Every number here is taken from the verified live deployment, not estimated.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "docs" / "Autonomous-Threat-Hunter.pptx"

# ---- FILL THESE IN -------------------------------------------------------
# The guidelines require team name, leader and members to be stated clearly.
TEAM_NAME = "CodeManiacs"
TEAM_LEADER = "Vaidehi Turkar"
TEAM_MEMBERS = "Vaidehi Turkar  ·  Shashwat"
# --------------------------------------------------------------------------

TRACK = "Domain 2: Cybersecurity  ·  Problem Statement 1"

# Dashboard palette (dashboard/src/index.css + severity.js) so the deck and the
# product read as one thing.
BG = RGBColor(0x0B, 0x0F, 0x14)
PANEL = RGBColor(0x15, 0x1B, 0x24)
BORDER = RGBColor(0x23, 0x2B, 0x36)
TEXT = RGBColor(0xE6, 0xE9, 0xEE)
DIM = RGBColor(0x7C, 0x87, 0x98)
DIMMER = RGBColor(0x4D, 0x58, 0x67)
ACCENT = RGBColor(0x4C, 0x9B, 0xE8)

CRITICAL = RGBColor(0xF0, 0x43, 0x5E)
HIGH = RGBColor(0xF5, 0x94, 0x3A)
MEDIUM = RGBColor(0xE8, 0xC5, 0x47)
LOW = RGBColor(0x4C, 0x9B, 0xE8)

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.8)
BODY_W = W - 2 * MARGIN

FONT = "Segoe UI"
MONO = "Consolas"


def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def textbox(slide, left, top, width, height, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.paragraphs[0].alignment = align
    return frame


def write(frame, text, size, color=TEXT, bold=False, font=FONT,
          space_after=0, space_before=0, align=None):
    p = frame.paragraphs[0] if not frame.paragraphs[0].runs else frame.add_paragraph()
    if align is not None:
        p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return p


def heading(slide, title, kicker=None):
    top = Inches(0.5)
    if kicker:
        f = textbox(slide, MARGIN, top, BODY_W, Inches(0.3))
        write(f, kicker.upper(), 11, ACCENT, bold=True)
        top += Inches(0.36)
    f = textbox(slide, MARGIN, top, BODY_W, Inches(0.7))
    write(f, title, 30, TEXT, bold=True)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, top + Inches(0.62),
                                  Inches(1.3), Pt(3))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False
    return top + Inches(0.98)


def card(slide, left, top, width, height, fill=PANEL):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = BORDER
    box.line.width = Pt(1)
    box.shadow.inherit = False
    box.adjustments[0] = 0.05
    return box


def stat_card(slide, left, top, width, value, label, color=TEXT, height=Inches(1.3)):
    card(slide, left, top, width, height)
    f = textbox(slide, left + Inches(0.24), top + Inches(0.2), width - Inches(0.48), Inches(0.55))
    write(f, value, 26, color, bold=True)
    f2 = textbox(slide, left + Inches(0.24), top + Inches(0.79), width - Inches(0.48), Inches(0.4))
    write(f2, label, 11, DIM)


def bullets(slide, top, items, size=15, gap=11, left=None, width=None):
    left = MARGIN if left is None else left
    width = BODY_W if width is None else width
    f = textbox(slide, left, top, width, H - top - Inches(0.5))
    for i, item in enumerate(items):
        lead, rest = item if isinstance(item, tuple) else (None, item)
        p = f.paragraphs[0] if i == 0 else f.add_paragraph()
        p.space_after = Pt(gap)
        dot = p.add_run()
        dot.text = "— "
        dot.font.size = Pt(size)
        dot.font.color.rgb = ACCENT
        dot.font.name = FONT
        if lead:
            r = p.add_run()
            r.text = lead + "  "
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = TEXT
            r.font.name = FONT
        r2 = p.add_run()
        r2.text = rest
        r2.font.size = Pt(size)
        r2.font.color.rgb = DIM if lead else TEXT
        r2.font.name = FONT
    return f


def footer(slide, text):
    f = textbox(slide, MARGIN, H - Inches(0.5), BODY_W, Inches(0.3))
    write(f, text, 9, DIMMER)


# ------------------------------------------------------------------ slides


def s1_title(prs):
    s = blank(prs)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.55), Pt(5), Inches(2.1))
    band.fill.solid(); band.fill.fore_color.rgb = ACCENT
    band.line.fill.background(); band.shadow.inherit = False

    f = textbox(s, MARGIN, Inches(1.5), BODY_W, Inches(0.35))
    write(f, TRACK.upper(), 12, ACCENT, bold=True)

    f = textbox(s, MARGIN, Inches(1.95), Inches(11.2), Inches(1.1))
    write(f, "Autonomous Threat Hunter", 46, TEXT, bold=True)

    f = textbox(s, MARGIN, Inches(2.95), Inches(11.2), Inches(0.8))
    write(f, "Insider attacks surfaced from raw enterprise logs — "
             "without a black box a SOC team can't question.", 17, DIM)

    # Team block: the guidelines require these stated clearly.
    top = Inches(4.15)
    card(s, MARGIN, top, Inches(6.0), Inches(1.85))
    f = textbox(s, MARGIN + Inches(0.3), top + Inches(0.22), Inches(5.4), Inches(1.4))
    write(f, TEAM_NAME, 19, TEXT, bold=True)
    write(f, "Team Leader   " + TEAM_LEADER, 12, DIM, font=MONO, space_before=10)
    write(f, "Members       " + TEAM_MEMBERS, 12, DIM, font=MONO, space_before=4)

    left2 = MARGIN + Inches(6.3)
    card(s, left2, top, BODY_W - Inches(6.3), Inches(1.85))
    f = textbox(s, left2 + Inches(0.3), top + Inches(0.22), Inches(5.0), Inches(1.4))
    write(f, "Live now", 13, ACCENT, bold=True)
    write(f, "threat-hunter-dashboard.onrender.com", 12, TEXT, font=MONO, space_before=9)
    write(f, "threat-hunter-api.onrender.com", 12, DIM, font=MONO, space_before=4)
    write(f, "github.com/Vaidehi2502/autonomous-threat-hunter", 12, DIM,
          font=MONO, space_before=4)


def s2_problem(prs):
    s = blank(prs)
    top = heading(s, "Insiders don't break in. They log in.", kicker="The problem")
    bullets(s, top, [
        ("No malware, no intrusion.", "An insider already holds valid credentials. "
         "Perimeter tools, antivirus and firewalls see nothing wrong."),
        ("The signal is purely behavioural.", "The only evidence is a change in how "
         "someone acts — hours kept, machines used, data moved."),
        ("The volume defeats humans.", "4,000 employees over 18 months is 1,393,138 "
         "user-days. Nobody reads that."),
        ("And the usual fix backfires.", "Drop in a black-box anomaly model and a SOC "
         "team gets alerts it cannot explain, cannot triage, and stops trusting. "
         "That is the false-positive trap this problem statement calls out."),
    ], size=16, gap=14)

    box_top = top + Inches(3.15)
    card(s, MARGIN, box_top, BODY_W, Inches(1.0), fill=RGBColor(0x11, 0x16, 0x1D))
    f = textbox(s, MARGIN + Inches(0.35), box_top + Inches(0.2), BODY_W - Inches(0.7), Inches(0.7))
    write(f, "Our answer, in one line", 11, DIMMER, bold=True)
    write(f, "Score every user-day against that person's own history — and make every "
             "flag explain itself.", 16, TEXT, space_before=6)
    footer(s, "Dataset: CERT r4.2 Insider Threat Test Dataset (Carnegie Mellon SEI) — synthetic, 4,000 users, 2010-2011")


def s3_solution(prs):
    s = blank(prs)
    top = heading(s, "Three decisions that make it trustworthy", kicker="Solution & innovation")

    gap = Inches(0.22)
    cw = (BODY_W - 2 * gap) / 3
    cards = [
        ("01", "Per-user baselines,", "not a population average",
         "A stockroom clerk and a physicist have different normals. Against a "
         "company-wide mean, one drowns real anomalies and the other hides them.",
         "Every user is compared only to themselves."),
        ("02", "ML is a second opinion,", "never the verdict",
         "An Isolation Forest catches multivariate patterns rules miss — but it "
         "only adds weight, never decides alone.",
         "A test asserts the model on its own can never reach High or Critical."),
        ("03", "Explains itself,", "two different ways",
         "Every reason names the exact counterfactual threshold. A separate "
         "14-day trend signal also catches sustained escalation no single day's "
         "z-score can see — without touching a single severity figure.",
         "Two failure modes covered, not one — see Slide 6 for the count."),
    ]
    for i, (num, line1, line2, body, footer_line) in enumerate(cards):
        left = MARGIN + i * (cw + gap)
        card(s, left, top, cw, Inches(3.05))
        f = textbox(s, left + Inches(0.24), top + Inches(0.22), cw - Inches(0.48), Inches(2.65))
        write(f, num, 12, ACCENT, bold=True, font=MONO)
        write(f, line1, 15, TEXT, bold=True, space_before=5)
        write(f, line2, 15, TEXT, bold=True)
        write(f, body, 11.5, DIM, space_before=10)
        write(f, footer_line, 11.5, TEXT, space_before=8)

    q_top = top + Inches(3.3)
    card(s, MARGIN, q_top, BODY_W, Inches(1.05), fill=RGBColor(0x11, 0x16, 0x1D))
    f = textbox(s, MARGIN + Inches(0.35), q_top + Inches(0.2), BODY_W - Inches(0.7), Inches(0.7))
    write(f, "Every flagged day carries a reason an analyst can act on or dismiss:",
          11, DIMMER, bold=True)
    write(f, '"off hours usb count unusually high (z=9.5) — would not flag below 1.1 (actual: 4)"',
          14, TEXT, font=MONO, space_before=7)


def s4_architecture(prs):
    s = blank(prs)
    top = heading(s, "Architecture & scoring engine", kicker="Technical implementation")

    stages = [("Raw logs", "logon · device\nusers"), ("Clean", "load_data.py\n→ parquet"),
              ("Features", "features.py\nper user-day"), ("Detect", "detect.py\nz-scores + ML"),
              ("Serve", "api.py\nFastAPI"), ("Analyse", "React\ndashboard")]
    gap = Inches(0.14)
    cw = (BODY_W - gap * (len(stages) - 1)) / len(stages)
    for i, (name, detail) in enumerate(stages):
        left = MARGIN + i * (cw + gap)
        hot = name in ("Detect", "Analyse")
        card(s, left, top, cw, Inches(1.5))
        f = textbox(s, left + Inches(0.1), top + Inches(0.2), cw - Inches(0.2),
                    Inches(0.3), align=PP_ALIGN.CENTER)
        write(f, name, 13, ACCENT if hot else TEXT, bold=True, align=PP_ALIGN.CENTER)
        f2 = textbox(s, left + Inches(0.08), top + Inches(0.63), cw - Inches(0.16),
                     Inches(0.8), align=PP_ALIGN.CENTER)
        for j, line in enumerate(detail.split("\n")):
            write(f2, line, 10, DIM, font=MONO, align=PP_ALIGN.CENTER,
                  space_before=0 if j == 0 else 2)
        if i < len(stages) - 1:
            a = textbox(s, left + cw, top + Inches(0.5), gap, Inches(0.35), align=PP_ALIGN.CENTER)
            write(a, "›", 17, DIMMER, align=PP_ALIGN.CENTER)

    mid = top + Inches(1.85)
    cw2 = (BODY_W - Inches(0.3)) / 2
    card(s, MARGIN, mid, cw2, Inches(2.35))
    f = textbox(s, MARGIN + Inches(0.28), mid + Inches(0.22), cw2 - Inches(0.56), Inches(1.9))
    write(f, "Five features, weighted", 15, ACCENT, bold=True)
    for feat, wt in [("off-hours logons", "3"), ("off-hours USB", "3"), ("USB connects", "2"),
                     ("distinct PCs", "2"), ("logon count", "1")]:
        write(f, f"{feat}  ·  weight {wt}", 12, DIM, font=MONO, space_before=6)

    left2 = MARGIN + cw2 + Inches(0.3)
    card(s, left2, mid, cw2, Inches(2.35))
    f = textbox(s, left2 + Inches(0.28), mid + Inches(0.22), cw2 - Inches(0.56), Inches(1.9))
    write(f, "Severity from the total", 15, ACCENT, bold=True)
    for lab, rng, col in [("Critical", "≥ 8", CRITICAL), ("High", "≥ 5", HIGH),
                          ("Medium", "≥ 2", MEDIUM), ("Low", "> 0", LOW)]:
        p = f.add_paragraph(); p.space_before = Pt(6)
        r = p.add_run(); r.text = f"{lab:<9}"; r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = col; r.font.name = MONO
        r2 = p.add_run(); r2.text = "score " + rng; r2.font.size = Pt(12)
        r2.font.color.rgb = DIM; r2.font.name = MONO
    write(f, "Rule weights past z = 2.5, plus 2 if the model agrees.", 12, DIM, space_before=9)

    footer(s, "Python · pandas · scikit-learn · PyArrow · FastAPI · React · Vite · Recharts · Docker · Render · GitHub Actions")


def s5_results(prs):
    s = blank(prs)
    top = heading(s, "What it found", kicker="Results")

    gap = Inches(0.24)
    cw = (BODY_W - 3 * gap) / 4
    for i, (val, lab, col) in enumerate([
        ("1,393,138", "user-days analysed", TEXT),
        ("4,000", "employees, 18 months", TEXT),
        ("50,520", "flagged  (3.6%)", ACCENT),
        ("377", "Critical days", CRITICAL),
    ]):
        stat_card(s, MARGIN + i * (cw + gap), top, cw, val, lab, col)

    bar_top = top + Inches(1.62)
    f = textbox(s, MARGIN, bar_top, Inches(4), Inches(0.25))
    write(f, "SEVERITY BREAKDOWN", 10, DIMMER, bold=True)
    row = bar_top + Inches(0.38)
    rows = [("Critical", 377, CRITICAL), ("High", 1641, HIGH),
            ("Medium", 47600, MEDIUM), ("Low", 902, LOW)]
    for label, value, col in rows:
        f = textbox(s, MARGIN, row, Inches(1.0), Inches(0.26))
        write(f, label, 12, col, bold=True)
        track = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN + Inches(1.1),
                                   row + Inches(0.04), Inches(4.3), Inches(0.18))
        track.fill.solid(); track.fill.fore_color.rgb = PANEL
        track.line.fill.background(); track.shadow.inherit = False
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN + Inches(1.1),
                                 row + Inches(0.04),
                                 max(Inches(0.05), Inches(4.3) * (value / 47600)), Inches(0.18))
        bar.fill.solid(); bar.fill.fore_color.rgb = col
        bar.line.fill.background(); bar.shadow.inherit = False
        f2 = textbox(s, MARGIN + Inches(5.55), row, Inches(1.2), Inches(0.26))
        write(f2, f"{value:,}", 12, DIM, font=MONO)
        row += Inches(0.42)

    left2 = MARGIN + Inches(7.1)
    card(s, left2, bar_top, BODY_W - Inches(7.1), Inches(2.4))
    f = textbox(s, left2 + Inches(0.3), bar_top + Inches(0.22),
                BODY_W - Inches(7.7), Inches(2.0))
    write(f, "Case study — EYD2871", 15, ACCENT, bold=True)
    write(f, "ProductionLineWorker, Assembly", 12, TEXT, space_before=7)
    write(f, "52 of 358 days flagged — 9 Critical, 21 High, 22 Medium — recurring "
             "across 13 months, not one bad week.", 12, DIM, space_before=6)
    write(f, "Same signature every time: off-hours logons, off-hours USB, and far more "
             "distinct PCs than a fixed workstation role should need.", 12, DIM,
          space_before=6)

    footer(s, "Second case ACM2278 (Salesman): two Critical days 12 apart, both 11.0, off-hours USB z=9.5 — exfiltration before departure. Every figure verified against the live API.")


def s6_engineering(prs):
    s = blank(prs)
    top = heading(s, "Built to be run, not just demoed", kicker="Engineering, security & deployment")

    gap = Inches(0.24)
    cw = (BODY_W - 3 * gap) / 4
    for i, (val, lab, col) in enumerate([
        ("80", "automated tests, CI on every push", ACCENT),
        ("330 MB", "memory, down from 870 MB", TEXT),
        ("< 1 s", "response over 1.39M rows", TEXT),
        ("Live", "deployed, not a localhost demo", ACCENT),
    ]):
        stat_card(s, MARGIN + i * (cw + gap), top, cw, val, lab, col)

    bullets(s, top + Inches(1.65), [
        ("Tests encode the argument.", "One asserts the same raw value is anomalous for "
         "a quiet user and normal for a busy one. Another asserts the ML model alone can "
         "never reach a High verdict, and a third proves a sustained trend is caught even "
         "when no single day's z-score is."),
        ("Measured, not guessed.", "Memory was profiled in a limited container — 256 MB "
         "is killed, 384 MB runs — which is why it fits a free instance."),
        ("Secured deliberately.", "Per-client rate limiting, named per-analyst API keys, "
         "hardened response headers, explicit CORS allowlist. Every read of an employee's "
         "timeline is audit-logged and attributed to the analyst who looked it up."),
    ], size=14, gap=11)

    link_top = top + Inches(3.75)
    card(s, MARGIN, link_top, BODY_W, Inches(0.72), fill=RGBColor(0x11, 0x16, 0x1D))
    f = textbox(s, MARGIN + Inches(0.32), link_top + Inches(0.19), BODY_W - Inches(0.64), Inches(0.4))
    write(f, "Try it   threat-hunter-dashboard.onrender.com      "
             "API   threat-hunter-api.onrender.com/docs", 13, ACCENT, font=MONO)


def s7_impact(prs):
    s = blank(prs)
    top = heading(s, "Where this goes", kicker="Scalability & real-world impact")

    cw = (BODY_W - Inches(0.3)) / 2
    card(s, MARGIN, top, cw, Inches(2.9))
    f = textbox(s, MARGIN + Inches(0.3), top + Inches(0.24), cw - Inches(0.6), Inches(2.4))
    write(f, "Scales without redesign", 16, ACCENT, bold=True)
    for line in ["Read-only and stateless — add instances, no coordination needed.",
                 "Detection runs offline and parallel; serving is a precomputed table.",
                 "Per-user baselines mean no retuning when new roles or teams appear.",
                 "Any log with a timestamp and an actor fits the same pipeline."]:
        write(f, "— " + line, 13, DIM, space_before=9)

    left2 = MARGIN + cw + Inches(0.3)
    card(s, left2, top, cw, Inches(2.9))
    f = textbox(s, left2 + Inches(0.3), top + Inches(0.24), cw - Inches(0.6), Inches(2.4))
    write(f, "Honest next steps", 16, ACCENT, bold=True)
    for line in ["47,600 Medium flags is a real triage load — tune thresholds, or treat "
                 "Medium as a watchlist rather than a page.",
                 "Only logon and device data so far. Email, file and web logs add the "
                 "most direct exfiltration signals.",
                 "Keys today are shared secrets in an env var — real rollout needs "
                 "self-service issuance, rotation, and per-analyst data scoping."]:
        write(f, "— " + line, 13, DIM, space_before=9)

    close_top = top + Inches(3.15)
    card(s, MARGIN, close_top, BODY_W, Inches(1.05), fill=RGBColor(0x11, 0x16, 0x1D))
    f = textbox(s, MARGIN + Inches(0.35), close_top + Inches(0.2), BODY_W - Inches(0.7), Inches(0.7))
    write(f, "Every alert explains itself.", 20, TEXT, bold=True)
    write(f, "That is what makes it usable by a real security team, instead of one more "
             "model they learn to ignore.", 14, DIM, space_before=6)


NOTES = [
    "We're {team}, on the Cybersecurity track, problem statement one — Autonomous "
    "Threat Hunter for Insider Attacks. It's built, deployed and live right now; "
    "the link is on screen and you can open it while I talk.",

    "Insider threats are the hardest category because nothing is technically broken. "
    "The person has valid credentials, so the only evidence is behavioural. Our "
    "dataset is 4,000 employees over 18 months — 1.39 million user-days. No human "
    "reads that. But the usual answer, a black-box anomaly model, gives a SOC team "
    "alerts it can't explain, so they stop trusting it. The problem statement "
    "explicitly asks for a low false-positive design philosophy, and that trap is "
    "exactly what we designed around.",

    "Three decisions. First, per-user baselines: a stockroom clerk and a physicist have "
    "totally different normals, so against a company average one drowns in noise and "
    "the other hides. Every person is compared only to themselves. Second, the ML is "
    "a second opinion, never the verdict — the Isolation Forest adds weight but can't "
    "decide alone. That isn't a promise, it's enforced in code and there's a test "
    "asserting the model on its own can never produce a High or Critical. Third, every "
    "flag explains itself two ways: the reason names the exact counterfactual threshold "
    "it crossed, like the one on screen, and a separate rolling-trend signal catches "
    "someone escalating slowly over weeks — never spiking hard enough on any single day "
    "to trip a z-score, which the primary system is structurally blind to. That signal "
    "doesn't touch severity at all, so every number on the results slide stays exact.",

    "The pipeline: raw CSVs, cleaned to parquet, features per user-day, detection, "
    "served over FastAPI to a React dashboard. Detection is precomputed, so the API "
    "just serves a scored table — that's how we get sub-second responses over 1.39 "
    "million rows. Five features, weighted by how suspicious they are, each converted "
    "to a z-score against that user's own history. Past 2.5 standard deviations the "
    "rule fires; add the weights, plus two if the model agrees, and that maps to a "
    "severity. Note the explainable rules outweigh the model.",

    "50,520 flagged out of 1.39 million — 3.6%, so it's selective, not crying wolf. "
    "377 Critical days. The case on the right is what convinced us it works: a "
    "production line worker, 52 of 358 days flagged, spread over 13 months. That's a "
    "pattern, not a bad week — and always the same signature. A second user, a "
    "salesman, shows a completely different shape: two Critical days twelve days "
    "apart, off-hours USB at nine standard deviations above his own norm. That's what "
    "data theft before resignation looks like.",

    "This was built to run, not just to demo. 80 tests, green on every push. Several of "
    "them encode the actual argument rather than just coverage. On memory — it "
    "started at 870 MB, which fits no free tier; we profiled it in a limited "
    "container, found the real bottleneck, and got it to 330 MB. And security: rate "
    "limiting, named per-analyst API keys, hardened headers, and every employee "
    "timeline lookup is audit-logged to the analyst who made it. Please open the link.",

    "On scale: read-only and stateless, so you add instances with no coordination, and "
    "per-user baselines mean no retuning as the org changes. Being honest about what's "
    "next — 47,600 Medium flags needs threshold tuning or a watchlist tier; we're only "
    "using logon and device data, and email and file logs would add the most direct "
    "exfiltration signals; and today's keys are shared secrets in an env var, so real "
    "rollout needs proper issuance and per-analyst data scoping. The thing to remember: "
    "every alert explains itself. "
    "That's what makes it usable instead of ignorable. Thank you.",
]


def build():
    prs = new_deck()
    makers = [s1_title, s2_problem, s3_solution, s4_architecture,
              s5_results, s6_engineering, s7_impact]
    assert len(makers) == len(NOTES) == 7, "guidelines require 6-7 slides"

    for maker, note in zip(makers, NOTES):
        maker(prs)
        slide = prs.slides[len(prs.slides._sldIdLst) - 1]
        slide.notes_slide.notes_text_frame.text = note.replace("{team}", TEAM_NAME)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}  ({len(prs.slides._sldIdLst)} slides, {OUT.stat().st_size / 1024:.0f} KB)")
    if "TEAM NAME" in TEAM_NAME:
        print("!! Team name/leader/members are still placeholders - edit the top of this file.")


if __name__ == "__main__":
    build()
